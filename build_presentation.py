"""
Build a 10-minute presentation deck (PPTX) for the Electricity Price Forecasting
project — content, bullet points, speaker notes, and key figures embedded.
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pathlib import Path

ROOT = Path(__file__).parent
FIG = ROOT / "artifacts" / "figures"
OUT = ROOT / "Electricity_Forecasting_Presentation.pptx"

NAVY = RGBColor(0x1F, 0x3B, 0x66)
GREEN = RGBColor(0x2C, 0xA0, 0x2C)
GREY = RGBColor(0x55, 0x55, 0x55)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

prs = Presentation()
prs.slide_width = Inches(13.333)   # 16:9
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]


def add_slide():
    return prs.slides.add_slide(BLANK)


def bg_bar(slide, color=NAVY, h=Inches(0.18)):
    bar = slide.shapes.add_shape(1, 0, 0, SW, h)
    bar.fill.solid(); bar.fill.fore_color.rgb = color
    bar.line.fill.background()


def title_box(slide, text, top=Inches(0.45), size=30, color=NAVY):
    tb = slide.shapes.add_textbox(Inches(0.6), top, SW - Inches(1.2), Inches(1.0))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = text
    p.font.size = Pt(size); p.font.bold = True; p.font.color.rgb = color
    return tb


def bullets(slide, items, top=Inches(1.7), left=Inches(0.7), width=None,
            height=None, size=18, gap=8):
    width = width or (SW - Inches(1.4))
    height = height or (SH - top - Inches(0.5))
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame; tf.word_wrap = True
    for i, it in enumerate(items):
        lvl = 0; txt = it
        if isinstance(it, tuple):
            lvl, txt = it
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.level = lvl
        p.space_after = Pt(gap)
        run = p.add_run(); run.text = ("• " if lvl == 0 else "– ") + txt
        run.font.size = Pt(size - 2 * lvl)
        run.font.color.rgb = RGBColor(0x22, 0x22, 0x22) if lvl == 0 else GREY
    return tb


def notes(slide, text):
    slide.notes_slide.notes_text_frame.text = text


def add_image(slide, name, left, top, width=None, height=None):
    return slide.shapes.add_picture(str(FIG / name), left, top, width=width, height=height)


# ───────────────────────── 1. TITLE ─────────────────────────
s = add_slide(); bg_bar(s, NAVY, Inches(0.25))
tb = s.shapes.add_textbox(Inches(0.8), Inches(2.2), SW - Inches(1.6), Inches(2.2))
tf = tb.text_frame; tf.word_wrap = True
p = tf.paragraphs[0]; p.text = "Electricity Price Forecasting with"
p.font.size = Pt(34); p.font.bold = True; p.font.color.rgb = NAVY
p2 = tf.add_paragraph(); p2.text = "Heterogeneous Graph Neural Networks"
p2.font.size = Pt(34); p2.font.bold = True; p2.font.color.rgb = NAVY
sub = s.shapes.add_textbox(Inches(0.8), Inches(4.4), SW - Inches(1.6), Inches(1.5))
stf = sub.text_frame; stf.word_wrap = True
sp = stf.paragraphs[0]
sp.text = "Multi-area day-ahead forecasting for the Nordic power market  •  ST-HeteroSAGE"
sp.font.size = Pt(18); sp.font.color.rgb = GREY
sp2 = stf.add_paragraph()
sp2.text = "Mahesh Gautam  ·  Riya Pokharel  ·  Sristi Kulung Rai  ·  Subhash Kumar Mehta"
sp2.font.size = Pt(15); sp2.font.color.rgb = GREY
notes(s, "10-minute talk. Intro line: We forecast day-ahead electricity prices for "
         "Denmark by treating the market as a graph and learning on it with GNNs. "
         "Our best model, ST-HeteroSAGE, beats a strong XGBoost baseline by 26.5%.")

# ───────────────────────── 2. PROBLEM ─────────────────────────
s = add_slide(); bg_bar(s); title_box(s, "The Problem: Why Forecasting Prices Is Hard")
bullets(s, [
    "Day-ahead electricity prices are increasingly volatile",
    (1, "Wind & solar growth, cross-border market coupling, 2021–23 fuel shocks"),
    (1, "Hourly prices can swing by an order of magnitude in a single day"),
    (1, "Prices even go negative when renewables outpace demand"),
    "Accurate forecasts have real economic value",
    (1, "Battery & EV operators, day-ahead/intraday trading, demand response"),
    "Limitation of most methods: each price zone modelled as an isolated time series",
    (1, "Neighbouring markets enter (if at all) as a few flat columns — no structure"),
], top=Inches(1.6))
notes(s, "~45s. Sell the problem: volatility is the challenge; isolated-zone modelling "
         "is the gap. Set up the graph idea on the next slide.")

# ───────────────────────── 3. KEY IDEA ─────────────────────────
s = add_slide(); bg_bar(s); title_box(s, "Key Idea: The Market IS a Graph")
bullets(s, [
    "Danish zones DK1 & DK2 are NOT independent",
    (1, "Electrically interconnected with Germany and the Nordic system"),
    (1, "Prices propagate along interconnectors when lines are uncongested;"),
    (1, "zones decouple when a line saturates"),
    "So the market is naturally a graph that changes hour by hour:",
    (1, "Nodes  = price areas (zone–hour states)"),
    (1, "Edges  = transmission corridors + temporal links"),
    "Graph Neural Networks (GNNs) are built for exactly this kind of data",
    (1, "Message passing lets information flow between connected nodes"),
], top=Inches(1.55))
notes(s, "~45s. This is the core insight and motivation for using GNNs. "
         "Emphasise: structure that flat models cannot represent.")

# ───────────────────────── 4. RQ & OBJECTIVES ─────────────────────────
s = add_slide(); bg_bar(s); title_box(s, "Research Question & Objectives")
tb = s.shapes.add_textbox(Inches(0.7), Inches(1.5), SW - Inches(1.4), Inches(1.3))
tf = tb.text_frame; tf.word_wrap = True
p = tf.paragraphs[0]
p.text = ("How can a heterogeneous graph be constructed and integrated into GNNs to "
          "improve the accuracy, interpretability, and robustness of multi-area "
          "day-ahead price forecasting in the Nordic market?")
p.font.size = Pt(17); p.font.italic = True; p.font.color.rgb = NAVY
bullets(s, [
    "Construct a heterogeneous graph of the interconnected market",
    "Build a ladder of models: XGBoost → Homogeneous GNN → Heterogeneous GNN → ST-HeteroSAGE",
    "Compare fairly — identical data, features & evaluation",
    "Explain it (interpretability) and stress-test it (robustness)",
    "Deliver a reproducible, end-to-end pipeline",
], top=Inches(3.1))
notes(s, "~40s. Read the RQ once. Then list objectives quickly — they map 1:1 to the "
         "rest of the talk.")

# ───────────────────────── 5. DATA ─────────────────────────
s = add_slide(); bg_bar(s); title_box(s, "Data")
bullets(s, [
    "5.75 years of hourly data: 2019-12-31 → 2025-09-30 (50,399 hours/zone)",
    "Four bidding zones:",
    (1, "DK1, DK2 (West/East Denmark) — forecast TARGETS — from Energinet/Nord Pool"),
    (1, "DE-LU (Germany), SE3 / 'HYDRO' (Sweden) — neighbour CONTEXT — from Energy-Charts"),
    "Features per node (17):",
    (1, "Price lags (24/48/168h) + 24h rolling mean/std"),
    (1, "Weather (temp, wind, cloud, humidity) via Open-Meteo"),
    (1, "Fundamentals (load, renewables, gas price, CO₂ price) + calendar (sin/cos)"),
    "Strict leakage control: chronological 80/10/10 split, scalers fit on train only",
], top=Inches(1.5), size=17)
notes(s, "~50s. Stress the leakage control — it's a key credibility point. "
         "DK1/DK2 are scored; DE & HYDRO are context.")

# ───────────────────────── 6. HETERO GRAPH ─────────────────────────
s = add_slide(); bg_bar(s); title_box(s, "The Heterogeneous Graph")
bullets(s, [
    "2 node types:",
    (1, "hour — 201,596 zone–hour states  |  market — 4 persistent zone nodes"),
    "5 edge types (typed relationships):",
    (1, "lag_to — temporal links (24/48/168h)"),
    (1, "co_occurs_with — same-hour cross-zone price co-movement"),
    (1, "belongs_to / rev — hour ↔ its market node"),
    (1, "interconnects — market ↔ market, weighted by MW transmission capacity"),
    "Distinctive: capacity-weighted interconnection — the model learns HOW STRONGLY zones are linked",
], top=Inches(1.55), size=17)
notes(s, "~45s. The capacity weighting is the novel modelling detail. Mention DE is "
         "excluded from co_occurs_with edges (stale 2024 data) — honest design choice.")

# ───────────────────────── 7. MODEL LADDER ─────────────────────────
s = add_slide(); bg_bar(s); title_box(s, "The Model Ladder — Each Fixes the Last")
bullets(s, [
    "XGBoost (baseline) — strong, but treats rows independently; no structure",
    (1, "MAE 205.6 / R² 0.52"),
    "Homogeneous GNN (GraphSAGE) — adds spatial+temporal edges → biggest jump",
    (1, "MAE 161.9 / R² 0.67  — but all nodes/edges are one type"),
    "GAT — tried attention; underperformed (honest negative result)  MAE 179.2",
    "HeteroSAGE — typed nodes/edges + capacity weights  MAE 162.8",
    (1, "Typing alone barely helped — temporal info still just 3 lag edges"),
    "ST-HeteroSAGE (proposed) — replaces lag edges with a causal temporal convolution",
    (1, "BEST: MAE 151.1 / R² 0.696"),
], top=Inches(1.5), size=16)
notes(s, "~60s. This is the narrative spine. Walk down the ladder: each model solves "
         "the previous one's weakness. End on ST-HeteroSAGE.")

# ───────────────────────── 8. ST-HETEROSAGE ─────────────────────────
s = add_slide(); bg_bar(s); title_box(s, "ST-HeteroSAGE Architecture")
bullets(s, [
    "Stack of 2 spatio-temporal blocks: SPATIAL message passing → TEMPORAL convolution",
    "Spatial: heterogeneous GraphSAGE across zones & market nodes (same hour)",
    "Temporal: Causal TCN (dilated 1-D convolutions) applied per zone",
    (1, "Kernel 7, dilations (1,4,24) → ~174h (≈ 1 week) receptive field"),
    (1, "CAUSAL = only looks at the past → no look-ahead leakage, by construction"),
    (1, "Replaces 3 discrete lag edges with a rich, continuous weekly view"),
    "BatchNorm + residuals, NO dropout (the two interact badly)",
    "Single shared regression head; loss masked to DK1 + DK2",
], top=Inches(1.55), size=16)
notes(s, "~55s. 'Causal' is the exam-favourite term — define it: only looks backward, "
         "guarantees no leakage. The ~1-week receptive field is why it wins.")

# ───────────────────────── 9. RESULTS ─────────────────────────
s = add_slide(); bg_bar(s); title_box(s, "Results — ST-HeteroSAGE Wins on Every Metric")
add_image(s, "02_model_comparison.png", Inches(0.5), Inches(1.5), width=Inches(7.6))
bullets(s, [
    "26.5% lower MAE than XGBoost",
    (1, "205.6 → 151.1 DKK/MWh"),
    "31% lower RMSE",
    "R² 0.52 → 0.70",
    "6.7% better than next-best graph model",
    "Same data & features for all —",
    (1, "the gain is the STRUCTURE, not more data"),
], left=Inches(8.3), top=Inches(1.7), width=Inches(4.7), size=15)
notes(s, "~55s. Headline slide. Two-step story: baseline→graph, then graph→temporal. "
         "Stress fairness: XGBoost had the same features incl. neighbour prices.")

# ───────────────────────── 10. ABLATION ─────────────────────────
s = add_slide(); bg_bar(s); title_box(s, "Ablation — What Actually Matters")
add_image(s, "05_ablation.png", Inches(0.5), Inches(1.6), width=Inches(8.0))
bullets(s, [
    "Strongest single finding:",
    "Remove SPATIAL → R² goes NEGATIVE (−0.26)",
    (1, "worse than predicting the mean"),
    "Remove TEMPORAL (TCN) → MAE +70%",
    "Both halves of 'spatio-temporal' are essential",
    "Market & hydro components: minor",
], left=Inches(8.7), top=Inches(1.7), width=Inches(4.3), size=15)
notes(s, "~45s. This proves the architecture earns its complexity. Spatial and "
         "temporal are both load-bearing — neither is decoration.")

# ───────────────────────── 11. INTERPRET + ROBUST ─────────────────────────
s = add_slide(); bg_bar(s); title_box(s, "Interpretability & Robustness")
bullets(s, [
    "Interpretability — the model recovers real market economics:",
    (1, "Top driver = natural-gas price (gas plants set the marginal price)"),
    (1, "Then price lags, renewables, load, carbon price"),
    (1, "Errors peak at the 17–18h evening peak and on Mondays — intuitive"),
    "Robustness — stress-tested at evaluation time:",
    (1, "Noise (+6.5% at 30%) and price spikes (+8% at 5×) → very robust"),
    (1, "Missing features (+30% at 30% dropout) → the one real weakness"),
    (1, "Implication: monitor data completeness, not outliers"),
], top=Inches(1.55), size=17)
notes(s, "~50s. Interpretability = trust (it learned economics, not just "
         "autocorrelation). Robustness = deployability (graceful under noise).")

# ───────────────────────── 12. ENGINEERING / DEMO ─────────────────────────
s = add_slide(); bg_bar(s); title_box(s, "Engineering & Live Demo")
bullets(s, [
    "Reproducible end-to-end pipeline (MLOps):",
    (1, "Ingest → SQLite → feature engineering → graph → train → serve"),
    (1, "Dockerised: MLflow tracking + FastAPI prediction service + Streamlit dashboard"),
    "Interactive dashboard (7 pages):",
    (1, "Live prediction, graph structure, forecast analysis, interpretability, robustness"),
    "Decision-support output: identifies the cheapest consumption hours day-ahead",
    "Deterministic & seeded for reproducibility (seed = 42)",
], top=Inches(1.6), size=17)
notes(s, "~45s. Show the dashboard live if possible (run LOCALLY — cloud may be flaky). "
         "Pick DK1, hour 18 → ~1063 DKK (evening peak). Fall back to figures if needed.")

# ───────────────────────── 13. CONCLUSION ─────────────────────────
s = add_slide(); bg_bar(s, GREEN); title_box(s, "Conclusion", color=GREEN)
bullets(s, [
    "Casting Nordic EPF as a heterogeneous spatio-temporal graph works:",
    (1, "Accuracy: best model, −26.5% MAE vs XGBoost (151.1 DKK, R² 0.70)"),
    (1, "Architecture: ablation proves spatial AND temporal are both essential"),
    (1, "Interpretability: recovers genuine market drivers (gas price dominant)"),
    (1, "Robustness: graceful under noise; sensitive only to missing data"),
    "Limitations: DE/HYDRO de-weighted (stale data); weather is a placeholder; single seed",
    "Future work: live weather feed, multi-seed runs, ARIMA/LSTM baselines, full DE/SE forecasting",
], top=Inches(1.55), size=17)
notes(s, "~45s. Close strong: the heterogeneous spatio-temporal graph is not just more "
         "complex — it's more accurate, explainable, and deployable. Be honest about "
         "limitations. Thank the panel and invite questions.")

# ───────────────────────── 14. THANK YOU ─────────────────────────
s = add_slide(); bg_bar(s, GREEN, Inches(0.25))
tb = s.shapes.add_textbox(Inches(0.8), Inches(2.8), SW - Inches(1.6), Inches(1.5))
tf = tb.text_frame
p = tf.paragraphs[0]; p.text = "Thank you — Questions?"
p.font.size = Pt(40); p.font.bold = True; p.font.color.rgb = NAVY
p2 = tf.add_paragraph()
p2.text = "ST-HeteroSAGE  ·  MAE 151.1 DKK  ·  R² 0.696  ·  −26.5% vs XGBoost"
p2.font.size = Pt(18); p2.font.color.rgb = GREY
notes(s, "Likely Q&A: 'Why does typing barely help but temporal helps a lot?' (ablation); "
         "'What is causal?' (only looks at past, no leakage); 'Why SMAPE not MAPE?' "
         "(prices near zero/negative break MAPE).")

prs.save(str(OUT))
print("Saved", OUT, "-", OUT.stat().st_size, "bytes,", len(prs.slides.__iter__.__self__._sldIdLst), "slides")
